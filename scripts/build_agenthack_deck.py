#!/usr/bin/env python3
"""Build the SWIMS-Connect AgentHack deck from the official UiPath PPTX template."""

from __future__ import annotations

import argparse
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


FONT = "Aptos"
WHITE = RGBColor(255, 255, 255)
INK = RGBColor(29, 38, 43)
MUTED = RGBColor(75, 91, 99)
TEAL = RGBColor(13, 160, 174)
DEEP_TEAL = RGBColor(0, 109, 101)
PALE_TEAL = RGBColor(224, 246, 246)
ORANGE = RGBColor(250, 70, 22)
GREEN = RGBColor(25, 166, 81)


def set_text(shape, text: str, size: float, *, color=INK, bold=False, align=None):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    for run in p.runs:
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tf


def set_bullets(shape, items: list[str], size: float, *, color=INK, spacing=8):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.05)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        p.line_spacing = 1.05
        # Suppress any bullet inherited by the template; the visible bullet is part of the text.
        p_pr = p._p.get_or_add_pPr()
        for tag in ("a:buChar", "a:buAutoNum", "a:buBlip", "a:buNone"):
            for node in p_pr.findall(tag, p_pr.nsmap):
                p_pr.remove(node)
        p_pr.append(OxmlElement("a:buNone"))
        for run in p.runs:
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return tf


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text,
    size,
    *,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    shape = slide.shapes.add_textbox(x, y, w, h)
    shape.text_frame.vertical_anchor = valign
    shape.text_frame.margin_left = Inches(0.08)
    shape.text_frame.margin_right = Inches(0.08)
    shape.text_frame.margin_top = Inches(0.04)
    shape.text_frame.margin_bottom = Inches(0.04)
    set_text(shape, text, size, color=color, bold=bold, align=align)
    return shape


def remove_shape(shape):
    element = shape._element
    element.getparent().remove(element)


def add_cropped_picture(slide, image_path: Path, x, y, w, h):
    with Image.open(image_path) as image:
        image_ratio = image.width / image.height
    box_ratio = w / h
    if image_ratio > box_ratio:
        picture = slide.shapes.add_picture(str(image_path), x, y, height=h)
        picture.left = x - int((picture.width - w) / 2)
    else:
        picture = slide.shapes.add_picture(str(image_path), x, y, width=w)
        picture.top = y - int((picture.height - h) / 2)

    # Crop overflow to the requested frame.
    if picture.width > w:
        overflow = (picture.width - w) / picture.width / 2
        picture.crop_left = overflow
        picture.crop_right = overflow
        picture.left = x
        picture.width = w
    if picture.height > h:
        overflow = (picture.height - h) / picture.height / 2
        picture.crop_top = overflow
        picture.crop_bottom = overflow
        picture.top = y
        picture.height = h
    return picture


def add_fitted_picture(slide, image_path: Path, x, y, w, h):
    """Fit an entire image inside a box without cropping and center it."""
    with Image.open(image_path) as image:
        image_ratio = image.width / image.height
    box_ratio = w / h
    if image_ratio > box_ratio:
        picture = slide.shapes.add_picture(str(image_path), x, y, width=w)
        picture.top = y + int((h - picture.height) / 2)
    else:
        picture = slide.shapes.add_picture(str(image_path), x, y, height=h)
        picture.left = x + int((w - picture.width) / 2)
    return picture


def style_node(shape, *, fill=PALE_TEAL, line=TEAL):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.5)
    shape.shadow.inherit = False


def add_node(slide, x, y, w, h, title, subtitle, *, fill=PALE_TEAL, line=TEAL):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    style_node(shape, fill=fill, line=line)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = title
    for run in p.runs:
        run.font.name = FONT
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = INK
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = subtitle
    p2.space_before = Pt(3)
    for run in p2.runs:
        run.font.name = FONT
        run.font.size = Pt(9.5)
        run.font.color.rgb = MUTED
    return shape


def add_arrow(slide, x1, y1, x2, y2, *, color=TEAL):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(2)
    line.line.end_arrowhead = True
    return line


def build(template: Path, output: Path, reporter_image: Path, worker_image: Path):
    prs = Presentation(str(template))
    if len(prs.slides) != 5:
        raise ValueError(f"Expected the official five-slide template, found {len(prs.slides)} slides")

    prs.core_properties.title = "SWIMS-Connect — UiPath AgentHack 2026"
    prs.core_properties.subject = "Conversational child-protection reporting and casework"
    prs.core_properties.author = "SWIMS-Connect"
    prs.core_properties.keywords = "UiPath, Maestro Case, Primero, child protection, WhatsApp"

    # Slide 1 — title.
    slide = prs.slides[0]
    set_text(slide.shapes[1], "SWIMS-Connect", 42, color=WHITE, bold=True)
    set_text(
        slide.shapes[0],
        "Conversational child-protection reporting and casework powered by UiPath",
        20,
        color=INK,
        bold=True,
    )
    slide.shapes[0].width = Inches(7.75)
    slide.shapes[0].height = Inches(1.15)
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            remove_shape(shape)
    add_fitted_picture(
        slide,
        reporter_image,
        Inches(9.05),
        Inches(1.1),
        Inches(2.45),
        Inches(4.15),
    )
    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.72),
        Inches(5.2),
        Inches(6.5),
        Inches(0.62),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = DEEP_TEAL
    badge.line.color.rgb = DEEP_TEAL
    set_text(
        badge,
        "UNICEF StartUp Lab Challenge Winner  •  Track 1: UiPath Maestro Case",
        13,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Slide 2 — problem and proposed solution.
    slide = prs.slides[1]
    set_text(slide.shapes[0], "Problem statement and proposed solution", 28, bold=True)
    set_text(slide.shapes[2], "Problem", 20, bold=True)
    set_bullets(
        slide.shapes[3],
        [
            "Child-protection concerns can be delayed by unfamiliar or inaccessible reporting channels.",
            "Frontline workers face repeated data entry, fragmented referral information and heavy caseloads.",
            "Missed assessments, referrals and follow-ups can delay support for children and families.",
            "Sensitive case information should not be copied into general-purpose AI tools.",
        ],
        15.5,
    )
    set_text(slide.shapes[4], "Solution", 20, bold=True)
    set_bullets(
        slide.shapes[5],
        [
            "Anonymous WhatsApp reporting through text, voice or images, creating a real Primero/SWIMS Case ID.",
            "Secure worker login with role-scoped case access, reports and casework assistance.",
            "Assessment, case-plan and referral drafts require explicit worker approval before saving.",
            "UiPath Maestro provides persistent deadline monitoring, verified against live Primero data.",
        ],
        15.5,
    )

    # Slide 3 — benefits and technologies.
    slide = prs.slides[2]
    set_text(slide.shapes[0], "Benefits and technologies used", 28, bold=True)
    table_shape = slide.shapes[2]
    table_shape.left = Inches(0.38)
    table_shape.top = Inches(1.35)
    table_shape.width = Inches(6.0)
    table_shape.height = Inches(5.1)
    table = table_shape.table
    rows = [
        ("End users", "Community reporters, social workers, supervisors"),
        ("User departments", "Government social welfare and child-protection agencies"),
        ("Industries", "Government, social services, humanitarian protection"),
        ("UiPath products used", "Automation Cloud, Coded Agents, Orchestrator, Maestro Case, TypeScript SDK"),
        ("Other integrations / APIs / technologies", "Python, LangGraph, Gemini, Node.js, WhatsApp, Primero REST API"),
    ]
    for row_index, (label, value) in enumerate(rows):
        for col_index, text in enumerate((label, value)):
            cell = table.cell(row_index, col_index)
            cell.text = text
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.04)
            cell.fill.solid()
            cell.fill.fore_color.rgb = TEAL
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = FONT
                    run.font.size = Pt(11.5 if row_index < 4 else 10.5)
                    run.font.bold = col_index == 0
                    run.font.color.rgb = WHITE if col_index == 0 else INK
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(3.5)
    slide.shapes[4].left = Inches(6.65)
    slide.shapes[4].top = Inches(1.25)
    slide.shapes[4].width = Inches(5.5)
    set_text(slide.shapes[4], "Benefits, impact and outcomes", 19, bold=True)
    slide.shapes[5].left = Inches(6.65)
    slide.shapes[5].top = Inches(1.75)
    slide.shapes[5].width = Inches(5.5)
    slide.shapes[5].height = Inches(4.9)
    set_bullets(
        slide.shapes[5],
        [
            "Familiar community access through WhatsApp, with anonymous reporting and consent controls.",
            "Less repetitive administration through structured intake, summaries, drafts and scheduled reports.",
            "Casework assistance without moving sensitive child information into an external AI chat.",
            "Proactive monitoring of assessments, plans, referrals, follow-ups and closure reviews.",
            "Configurable for Primero programmes across more than 80 countries and territories.",
        ],
        15,
        spacing=10,
    )

    # Slide 4 — architecture.
    slide = prs.slides[3]
    set_text(slide.shapes[2], "Solution architecture", 28, bold=True)
    remove_shape(slide.shapes[3])

    y = Inches(2.05)
    h = Inches(1.25)
    w = Inches(2.45)
    x_positions = [Inches(0.45), Inches(3.2), Inches(6.0), Inches(8.85)]
    nodes = [
        add_node(slide, x_positions[0], y, w, h, "Community & workers", "WhatsApp text • voice • images"),
        add_node(slide, x_positions[1], y, w, h, "WhatsApp gateway", "Consent • login • media • secure context"),
        add_node(
            slide,
            x_positions[2],
            y,
            w,
            h,
            "UiPath coded agent",
            "Python • LangGraph • Gemini • typed tools",
            fill=RGBColor(222, 241, 255),
            line=RGBColor(0, 120, 212),
        ),
        add_node(
            slide,
            x_positions[3],
            y,
            w,
            h,
            "Primero / SWIMS",
            "System of record • roles • workflow • audit",
            fill=RGBColor(235, 247, 238),
            line=GREEN,
        ),
    ]
    for left, right in zip(nodes, nodes[1:]):
        add_arrow(
            slide,
            left.left + left.width,
            left.top + int(left.height / 2),
            right.left,
            right.top + int(right.height / 2),
        )

    orchestrator = add_node(
        slide,
        Inches(3.2),
        Inches(4.3),
        Inches(3.65),
        Inches(1.15),
        "UiPath Orchestrator",
        "Packages • assets • secrets • jobs • traces",
        fill=RGBColor(245, 240, 255),
        line=RGBColor(110, 69, 226),
    )
    maestro = add_node(
        slide,
        Inches(7.15),
        Inches(4.3),
        Inches(3.75),
        Inches(1.15),
        "UiPath Maestro Case",
        "Persistent SLA clocks • verified overdue reminders",
        fill=RGBColor(255, 240, 229),
        line=ORANGE,
    )
    add_arrow(
        slide,
        nodes[2].left + int(nodes[2].width / 2),
        nodes[2].top + nodes[2].height,
        orchestrator.left + int(orchestrator.width / 2),
        orchestrator.top,
        color=RGBColor(110, 69, 226),
    )
    add_arrow(
        slide,
        nodes[3].left + int(nodes[3].width / 2),
        nodes[3].top + nodes[3].height,
        maestro.left + int(maestro.width / 2),
        maestro.top,
        color=ORANGE,
    )
    add_arrow(
        slide,
        maestro.left,
        maestro.top + int(maestro.height / 2),
        orchestrator.left + orchestrator.width,
        orchestrator.top + int(orchestrator.height / 2),
        color=ORANGE,
    )
    add_textbox(
        slide,
        Inches(0.5),
        Inches(6.05),
        Inches(11.6),
        Inches(0.38),
        "github.com/zachryit/swims-connect-uipath",
        13,
        color=DEEP_TEAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # Slide 5 — closing. Replace the source's special layout because its artwork overlays all
    # slide-level content; use the template's dark logo layout so our content remains editable.
    old_slide_id = prs.slides._sldIdLst[-1]
    prs.part.drop_rel(old_slide_id.rId)
    del prs.slides._sldIdLst[-1]
    slide = prs.slides.add_slide(prs.slide_layouts[19])
    add_textbox(
        slide,
        Inches(0.62),
        Inches(0.62),
        Inches(6.4),
        Inches(2.2),
        "THANK\nYOU.",
        58,
        color=ORANGE,
        bold=True,
    )
    add_textbox(
        slide,
        Inches(0.72),
        Inches(3.45),
        Inches(6.0),
        Inches(1.1),
        "Helping every child-protection concern reach the people who can act.",
        19,
        color=WHITE,
        bold=True,
    )
    add_textbox(
        slide,
        Inches(0.72),
        Inches(4.75),
        Inches(6.4),
        Inches(1.1),
        "WhatsApp  +233 54 159 9802\nGitHub  github.com/zachryit/swims-connect-uipath",
        14,
        color=WHITE,
    )
    add_fitted_picture(
        slide,
        reporter_image,
        Inches(7.25),
        Inches(0.7),
        Inches(2.25),
        Inches(5.6),
    )
    add_fitted_picture(
        slide,
        worker_image,
        Inches(9.85),
        Inches(0.7),
        Inches(2.25),
        Inches(5.6),
    )

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--template", type=Path, required=True)
    parser.add_argument("--output", type=Path, required=True)
    parser.add_argument(
        "--reporter-image",
        type=Path,
        default=Path("docs/images/swims-connect-conversation.png"),
    )
    parser.add_argument(
        "--worker-image",
        type=Path,
        default=Path("docs/images/swims-connect-worker-conversation.png"),
    )
    args = parser.parse_args()
    build(args.template, args.output, args.reporter_image, args.worker_image)


if __name__ == "__main__":
    main()
